import os
import sys
import zipfile
from PIL import Image, ImageSequence
from pptx import Presentation
from pptx.util import Inches

from PyQt6.QtWidgets import (QApplication, QWidget, QVBoxLayout, QHBoxLayout, 
                             QPushButton, QLineEdit, QLabel, QFileDialog, 
                             QMessageBox, QProgressBar)
from PyQt6.QtCore import Qt


def modify_gif_loop_count(gif_path):
    """ تقرأ الـ GIF وتُعيد حفظها لتشتغل الدورة الأولى فقط وتتوقف فوراً دون أي تكرار """
    try:
        with Image.open(gif_path) as img:
            if not img.is_animated:
                return
            
            # استخراج جميع فريمات الصورة مع المدة الزمنية لكل فريم
            frames = []
            durations = []
            for frame in ImageSequence.Iterator(img):
                frames.append(frame.copy())
                durations.append(frame.info.get('duration', 100))
            
            # إعادة حفظها بدون خاصية التكرار (loop) لتعمل مرة واحدة فقط وتتوقف عند الفريم الأخير
            if frames:
                frames[0].save(
                    gif_path,
                    save_all=True,
                    append_images=frames[1:],
                    duration=durations,
                    loop=None  # None تعني تشغيل الدورة الأصلية فقط دون أي إعادة (No repeat)
                )
    except Exception as e:
        print(f"فشل تعديل تكرار الصورة {gif_path}: {e}")


class GifToPptxApp(QWidget):
    def __init__(self):
        super().__init__()
        self.initUI()
        
    def initUI(self):
        # إعدادات النافذة الرئيسية
        self.setWindowTitle('مُحول GIF ZIP إلى بوربوينت')
        self.resize(550, 250)
        self.setLayoutDirection(Qt.LayoutDirection.RightToLeft) # واجهة عربية متناسقة
        
        # التصميم الرئيسي (عمودي)
        main_layout = QVBoxLayout()
        
        # --- اختيار ملف الـ ZIP ---
        zip_label = QLabel("اختر ملف الـ ZIP المحتوي على صور GIF:")
        main_layout.addWidget(zip_label)
        
        zip_h_layout = QHBoxLayout()
        self.zip_input = QLineEdit()
        self.zip_input.setPlaceholderText("مسار ملف الـ ZIP...")
        self.btn_browse_zip = QPushButton("استعراض...")
        self.btn_browse_zip.clicked.connect(self.browse_zip)
        zip_h_layout.addWidget(self.zip_input)
        zip_h_layout.addWidget(self.btn_browse_zip)
        main_layout.addLayout(zip_h_layout)
        
        # --- اختيار مسار حفظ البوربوينت ---
        save_label = QLabel("حدد مكان حفظ ملف البوربوينت الناتج:")
        main_layout.addWidget(save_label)
        
        save_h_layout = QHBoxLayout()
        self.save_input = QLineEdit()
        self.save_input.setPlaceholderText("مكان حفظ الملف ونوعه .pptx ...")
        self.btn_browse_save = QPushButton("تحديد الحفظ...")
        self.btn_browse_save.clicked.connect(self.browse_save)
        save_h_layout.addWidget(self.save_input)
        save_h_layout.addWidget(self.btn_browse_save)
        main_layout.addLayout(save_h_layout)
        
        # --- شريط التقدم ---
        self.progress_bar = QProgressBar()
        self.progress_bar.setValue(0)
        main_layout.addWidget(self.progress_bar)
        
        # --- زر البدء والتنفيذ ---
        self.btn_start = QPushButton("ابدأ عملية التحويل المباشر")
        self.btn_start.setStyleSheet("background-color: #2b78e4; color: white; font-weight: bold; font-size: 14px; padding: 8px;")
        self.btn_start.clicked.connect(self.process_conversion)
        main_layout.addWidget(self.btn_start)
        
        self.setLayout(main_layout)

    def browse_zip(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "اختر ملف ZIP", "", "Zip Files (*.zip)")
        if file_path:
            self.zip_input.setText(file_path)
            
    def browse_save(self):
        file_path, _ = QFileDialog.getSaveFileName(self, "حفظ ملف البوربوينت", "Presentation.pptx", "PowerPoint Files (*.pptx)")
        if file_path:
            if not file_path.lower().endswith('.pptx'):
                file_path += '.pptx'
            self.save_input.setText(file_path)

    def process_conversion(self):
        zip_path = self.zip_input.text().strip()
        output_pptx_path = self.save_input.text().strip()
        
        # التحقق من المدخلات
        if not zip_path or not os.path.exists(zip_path):
            QMessageBox.warning(self, "خطأ", "الرجاء اختيار ملف ZIP صحيح وموجود.")
            return
        if not output_pptx_path:
            QMessageBox.warning(self, "خطأ", "الرجاء تحديد مسار حفظ ملف البوربوينت.")
            return

        self.btn_start.setEnabled(False)
        self.progress_bar.setValue(10)
        
        # إعداد البوربوينت بأبعاد عريضة 16:9
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[6]
        
        extract_dir = "temp_ui_gifs"
        os.makedirs(extract_dir, exist_ok=True)
        
        try:
            # 1. فك الضغط
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                zip_ref.extractall(extract_dir)
            self.progress_bar.setValue(30)
            
            # تجميع الصور
            gif_files = []
            for root, dirs, files in os.walk(extract_dir):
                for file in files:
                    if file.lower().endswith('.gif'):
                        gif_files.append(os.path.join(root, file))
            
            gif_files.sort()
            
            if not gif_files:
                QMessageBox.information(self, "تنبيه", "لم يتم العثور على أي صور GIF داخل ملف الـ ZIP المختار.")
                self.reset_ui()
                return
            
            total_files = len(gif_files)
            
            # 2. تعديل خاصية الـ Loop وإدراجها بالسلايدات
            for index, gif_path in enumerate(gif_files):
                modify_gif_loop_count(gif_path)
                
                # إضافة شريحة
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # جعل الصورة تغطي كامل مساحة الشريحة
                slide.shapes.add_picture(gif_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                
                # تحديث شريط التقدم ديناميكياً
                progress_val = 30 + int((index + 1) / total_files * 60)
                self.progress_bar.setValue(progress_val)
                
            # 3. حفظ الملف النهائي
            prs.save(output_pptx_path)
            self.progress_bar.setValue(100)
            QMessageBox.information(self, "نجاح", f"تم تحويل الصور وصناعة العرض التقديمي بنجاح!\nالمسار: {output_pptx_path}")
            
        except Exception as e:
            QMessageBox.critical(self, "خطأ غير متوقع", f"فشلت العملية بسبب الخطأ التالي:\n{e}")
            
        finally:
            # تنظيف المجلد المؤقت وحذف الصور المستخرجة
            for root, dirs, files in os.walk(extract_dir, topdown=False):
                for file in files:
                    os.remove(os.path.join(root, file))
                for dir in dirs:
                    os.rmdir(os.path.join(root, dir))
            if os.path.exists(extract_dir):
                os.rmdir(extract_dir)
            
            self.reset_ui()

    def reset_ui(self):
        self.btn_start.setEnabled(True)
        self.progress_bar.setValue(0)


if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = GifToPptxApp()
    ex.show()
    sys.exit(app.exec())